#!/usr/bin/env python3
"""Generate a one-slide PowerPoint overview of Scout (16:9)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Scout brand palette (from tailwind.config.ts) ----
INK_950 = RGBColor(0x0A, 0x0A, 0x0A)
INK_900 = RGBColor(0x1D, 0x1D, 0x1F)
INK_700 = RGBColor(0x3A, 0x3A, 0x3C)
INK_300 = RGBColor(0xA1, 0xA1, 0xA6)
INK_100 = RGBColor(0xE8, 0xE8, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0x71, 0xE3)   # accent-500
FUCHSIA = RGBColor(0xD9, 0x46, 0xEF)
PINK = RGBColor(0xEC, 0x48, 0x99)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
CARD = RGBColor(0x16, 0x16, 0x18)
CARD_BORDER = RGBColor(0x33, 0x33, 0x38)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank


def set_gradient(shape, stops, angle=45):
    """Apply a linear gradient fill to a shape via raw XML."""
    spPr = shape.fill._xPr  # element holding fill props
    # remove existing fill
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill",
                "a:pattFill", "a:grpFill"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    grad = spPr.makeelement(qn("a:gradFill"), {})
    gsLst = grad.makeelement(qn("a:gsLst"), {})
    for pos, color in stops:
        gs = grad.makeelement(qn("a:gs"), {"pos": str(int(pos * 1000))})
        srgb = grad.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % (color[0], color[1], color[2])})
        gs.append(srgb)
        gsLst.append(gs)
    grad.append(gsLst)
    lin = grad.makeelement(qn("a:lin"), {"ang": str(int(angle * 60000)), "scaled": "1"})
    grad.append(lin)
    # insert gradient before a:ln if present
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        ln.addprevious(grad)
    else:
        spPr.append(grad)


def no_line(shape):
    shape.line.fill.background()


def add_rect(x, y, w, h, fill=None, line=None, line_w=None, radius=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        x, y, w, h)
    shp.shadow.inherit = False
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = line_w or Pt(1)
    else:
        no_line(shp)
    return shp


def add_text(x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=0, line_spacing=None):
    """runs: list of paragraphs; each paragraph is list of (text,size,color,bold,spacing)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after:
            p.space_after = Pt(space_after)
        if line_spacing:
            p.line_spacing = line_spacing
        for (text, size, color, bold, spacing) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Arial"
            if spacing is not None:
                rPr = r._r.get_or_add_rPr()
                rPr.set("spc", str(int(spacing)))
    return tb


# ---- Background ----
bg = add_rect(0, 0, SW, SH, fill=INK_950, radius=False)
set_gradient(bg, [(0, (0x0A, 0x0A, 0x0F)), (55, (0x12, 0x10, 0x1E)),
                  (100, (0x05, 0x05, 0x0A))], angle=115)

# Aurora glow accents (soft rounded blobs)
glow1 = add_rect(Inches(-1.2), Inches(-1.4), Inches(5.2), Inches(4.6), fill=ACCENT)
glow1.fill.fore_color.rgb = ACCENT
glow1.adjustments[0] = 0.5
set_gradient(glow1, [(0, (0x00, 0x71, 0xE3)), (100, (0x0A, 0x0A, 0x0F))], angle=45)
glow1.line.fill.background()

glow2 = add_rect(Inches(9.6), Inches(3.4), Inches(5.4), Inches(5.0), fill=PINK)
glow2.adjustments[0] = 0.5
set_gradient(glow2, [(0, (0xEC, 0x48, 0x99)), (60, (0x6E, 0x21, 0x6E)),
                     (100, (0x0A, 0x0A, 0x0F))], angle=300)
glow2.line.fill.background()

# ---- Eyebrow pill ----
pill = add_rect(Inches(0.85), Inches(0.62), Inches(3.65), Inches(0.42),
                fill=RGBColor(0x1A, 0x1A, 0x20), line=CARD_BORDER, line_w=Pt(0.75))
dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.02), Inches(0.755),
                             Inches(0.11), Inches(0.11))
dot.fill.solid(); dot.fill.fore_color.rgb = EMERALD; no_line(dot); dot.shadow.inherit = False
add_text(Inches(1.22), Inches(0.62), Inches(3.3), Inches(0.42),
         [[("BUILT FOR BROADBAND SERVICE PROVIDERS", 9.5, INK_300, True, 60)]],
         anchor=MSO_ANCHOR.MIDDLE)

# ---- Title ----
add_text(Inches(0.8), Inches(1.18), Inches(8.6), Inches(2.1),
         [[("Scout", 58, WHITE, True, -20)]], anchor=MSO_ANCHOR.TOP)
# Tagline (two-tone)
add_text(Inches(0.82), Inches(2.18), Inches(8.8), Inches(1.0),
         [[("Know your market ", 30, WHITE, True, -10),
           ("before they do.", 30, ACCENT, True, -10)]])

# ---- Subhead ----
add_text(Inches(0.85), Inches(3.02), Inches(7.7), Inches(1.3),
         [[("Scout turns a list of ZIP codes into a complete competitive "
            "briefing for broadband providers — overlap, technology mix, "
            "demographics, sentiment, and recent launches, plus a ranked list "
            "of opportunities for your next move.", 14.5, INK_100, False, 0)]],
         line_spacing=1.18)

# ---- Stat chips ----
stats = [("FCC BDC", "coverage data"), ("Census ACS", "demographics"),
         ("A–F", "AI readiness"), ("1 link", "shareable")]
sx = Inches(0.85)
chip_w = Inches(1.72)
gap = Inches(0.18)
for i, (big, small) in enumerate(stats):
    x = Emu(int(sx) + i * (int(chip_w) + int(gap)))
    c = add_rect(x, Inches(4.35), chip_w, Inches(0.92),
                 fill=CARD, line=CARD_BORDER, line_w=Pt(0.75))
    add_text(x, Inches(4.46), chip_w, Inches(0.45),
             [[(big, 16, WHITE, True, 0)]], align=PP_ALIGN.CENTER)
    add_text(x, Inches(4.86), chip_w, Inches(0.34),
             [[(small.upper(), 8, INK_300, True, 40)]], align=PP_ALIGN.CENTER)

# ---- "What it does" feature cards (right column) ----
features = [
    ("Competitor overlap", "Every provider in your footprint — technology, speeds, and where overlap is densest.", ACCENT),
    ("Ranked opportunities", "Whitespace, counter-launches & civic plays — each tied to evidence from your ZIPs.", FUCHSIA),
    ("AI readiness score", "Can AI assistants find & recommend you? Graded A–F on crawlers, schema & content.", PINK),
    ("Sentiment + launch radar", "Competitor review scores and what rivals just launched — FWA, fiber, smart home.", EMERALD),
]
fx = Inches(8.85)
fw = Inches(4.0)
fy = Inches(0.62)
fh = Inches(1.46)
fgap = Inches(0.13)
for i, (title, body, accent) in enumerate(features):
    y = Emu(int(fy) + i * (int(fh) + int(fgap)))
    card = add_rect(fx, y, fw, fh, fill=CARD, line=CARD_BORDER, line_w=Pt(0.75))
    # accent bar
    bar = add_rect(fx, y, Inches(0.09), fh, fill=accent, radius=False)
    add_text(Emu(int(fx) + Inches(0.28)), Emu(int(y) + Inches(0.18)),
             Emu(int(fw) - int(Inches(0.5))), Inches(0.4),
             [[(title, 14.5, WHITE, True, 0)]])
    add_text(Emu(int(fx) + Inches(0.28)), Emu(int(y) + Inches(0.62)),
             Emu(int(fw) - int(Inches(0.5))), Inches(0.8),
             [[(body, 10.5, INK_300, False, 0)]], line_spacing=1.12)

# ---- Footer ----
add_text(Inches(0.85), Inches(6.85), Inches(11.6), Inches(0.4),
         [[("No login  ·  Shareable link  ·  Next.js 15 on Vercel  ·  "
            "Solutions: SmartHome · SmartTown · SmartMDU · "
            "Home OfficeIQ · SmartBiz", 9.5, INK_300, False, 20)]])

out = "/home/user/scout/Scout-overview.pptx"
prs.save(out)
print("saved", out)
